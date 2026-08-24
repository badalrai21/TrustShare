import csv
import io
import re
from pathlib import Path

from fastapi import HTTPException, status

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".csv", ".xlsx"}


def _clean(text: str) -> str:
    text = text.replace("\x00", " ")
    return re.sub(r"[ \t]+", " ", re.sub(r"\n{3,}", "\n\n", text)).strip()


def extract_text(filename: str, data: bytes) -> str:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
            detail={"code": "unsupported_file_type", "message": f"Summarisation is unavailable for {extension or 'this file format'}."},
        )
    try:
        if extension in {".txt", ".md"}:
            text = data.decode("utf-8-sig")
        elif extension == ".pdf":
            from pypdf import PdfReader
            text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(data)).pages)
        elif extension == ".docx":
            from docx import Document
            text = "\n".join(p.text for p in Document(io.BytesIO(data)).paragraphs)
        elif extension == ".pptx":
            from pptx import Presentation
            text = "\n".join(shape.text for slide in Presentation(io.BytesIO(data)).slides for shape in slide.shapes if hasattr(shape, "text"))
        elif extension == ".csv":
            rows = csv.reader(io.StringIO(data.decode("utf-8-sig")))
            text = "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
        else:
            from openpyxl import load_workbook
            workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            text = "\n".join(
                " | ".join("" if value is None else str(value) for value in row)
                for sheet in workbook.worksheets for row in sheet.iter_rows(values_only=True)
            )
            workbook.close()
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=400, detail={"code": "corrupt_document", "message": "The document could not be read safely."}) from exc
    text = _clean(text)
    if not text:
        message = "This document does not contain extractable text. OCR support is not currently enabled." if extension == ".pdf" else "This document does not contain extractable text."
        raise HTTPException(status_code=400, detail={"code": "no_extractable_text", "message": message})
    return text


def chunk_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    if chunk_size <= overlap:
        raise ValueError("Chunk size must be greater than overlap")
    chunks, start = [], 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        if end < len(text):
            boundary = max(text.rfind(". ", start + chunk_size // 2, end), text.rfind("\n", start + chunk_size // 2, end))
            if boundary > start:
                end = boundary + 1
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start = end - overlap
    return [chunk for chunk in chunks if chunk]

